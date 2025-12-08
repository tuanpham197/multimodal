import base64
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.pptx import partition_pptx
from unstructured.documents.elements import Table, Title
from langchain_text_splitters import RecursiveCharacterTextSplitter


class PptxParser:
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self._text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ".", ":", ";", ",", " ", ""],
        )

    def extract_images(self, file_path: str) -> list[str]:
        prs = Presentation(file_path)
        images_base64 = []

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    image_base64 = base64.b64encode(image_bytes).decode("utf-8")
                    images_base64.append(image_base64)
                    print(f"[DEBUG PptxParser] Slide {slide_num}: Found image")

                if hasattr(shape, "shapes"):
                    for sub_shape in shape.shapes:
                        if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_bytes = sub_shape.image.blob
                            image_base64 = base64.b64encode(image_bytes).decode("utf-8")
                            images_base64.append(image_base64)
                            print(f"[DEBUG PptxParser] Slide {slide_num}: Found grouped image")

        print(f"[DEBUG PptxParser] Total images extracted: {len(images_base64)}")
        return images_base64

    def extract_text_and_tables(self, file_path: str) -> dict:
        elements = partition_pptx(filename=file_path)
        print(f"\n[DEBUG PptxParser] Raw elements: {len(elements)}")

        sections = self._group_by_sections(elements)
        print(f"[DEBUG PptxParser] Grouped into {len(sections)} sections")

        all_chunks = []
        tables = []

        for section in sections:
            if section["type"] == "table":
                tables.append(section["content"])
            else:
                section_text = section["content"]
                if len(section_text) > self._text_splitter._chunk_size:
                    chunks = self._text_splitter.split_text(section_text)
                    all_chunks.extend(chunks)
                elif len(section_text) >= 50:
                    all_chunks.append(section_text)

        print(f"\n[DEBUG PptxParser] Final chunks: {len(all_chunks)} texts, {len(tables)} tables")
        for i, chunk in enumerate(all_chunks[:5]):
            print(f"[DEBUG PptxParser] Chunk[{i}] ({len(chunk)} chars): {chunk[:150]}...")

        return {"texts": all_chunks, "tables": tables}

    def _group_by_sections(self, elements: list) -> list[dict]:
        sections = []
        current_section = []

        for element in elements:
            if isinstance(element, Table):
                if current_section:
                    sections.append({
                        "type": "text",
                        "content": self._merge_texts(current_section)
                    })
                    current_section = []
                table_text = element.metadata.text_as_html if hasattr(element.metadata, 'text_as_html') else str(element)
                sections.append({"type": "table", "content": table_text})
                continue

            text = element.text if hasattr(element, 'text') else str(element)
            if not text or not text.strip():
                continue

            is_major_title = isinstance(element, Title) and self._is_major_section(text)

            if is_major_title and current_section:
                sections.append({
                    "type": "text",
                    "content": self._merge_texts(current_section)
                })
                current_section = []

            current_section.append(text)

        if current_section:
            sections.append({
                "type": "text",
                "content": self._merge_texts(current_section)
            })

        return sections

    def _is_major_section(self, text: str) -> bool:
        text = text.strip()
        if re.match(r'^[IVXLCDM]+\.\s+', text):
            return True
        if re.match(r'^[A-Z]\.\s+', text):
            return True
        if re.match(r'^(Slide|Phần|Chương)\s+\d+', text, re.IGNORECASE):
            return True
        return False

    def _merge_texts(self, texts: list[str]) -> str:
        return "\n".join(t.strip() for t in texts if t.strip())

    def parse(self, file_path: str) -> dict:
        images = self.extract_images(file_path)
        elements = self.extract_text_and_tables(file_path)

        print(f"\n[DEBUG PptxParser] === PARSE RESULT ===")
        print(f"[DEBUG PptxParser] Images: {len(images)}")
        print(f"[DEBUG PptxParser] Texts: {len(elements['texts'])}")
        print(f"[DEBUG PptxParser] Tables: {len(elements['tables'])}")

        return {
            "images": images,
            "texts": elements["texts"],
            "tables": elements["tables"],
        }
