from unstructured.partition.pdf import partition_pdf
from unstructured.documents.elements import Table, Image, CompositeElement
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.embeddings import Embeddings
from langchain_groq import ChatGroq
from pathlib import Path
from typing import List
from langchain_chroma import Chroma
from langchain_ollama import ChatOllama
from langchain_core.stores import InMemoryStore
from langchain_core.documents import Document
from langchain_classic.retrievers.multi_vector import MultiVectorRetriever
from langchain_google_genai import ChatGoogleGenerativeAI

from langchain_core.runnables import RunnablePassthrough, RunnableLambda
from langchain_core.messages import SystemMessage, HumanMessage
from langchain_openai import ChatOpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from base64 import b64decode
from unstructured.partition.pptx import partition_pptx
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.auto import partition
from unstructured.documents.elements import (
    Element,
    NarrativeText,
    Title,
    Table,
    Image,
)


import chromadb
import ollama
import os
import uuid
import re
import base64
import docx2txt

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")



class OllamaEmbeddingsCustom(Embeddings):
    def __init__(self, model: str = "nomic-embed-text", max_length: int = 2000):
        self.model = model
        self.max_length = max_length
        self.client = ollama.Client(host="http://localhost:11434")

    def _truncate(self, text: str) -> str:
        if len(text) > self.max_length:
            print(f"[DEBUG] Truncating text from {len(text)} to {self.max_length} chars")
            return text[:self.max_length]
        return text

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        embeddings = []
        for i, text in enumerate(texts):
            truncated_text = self._truncate(text)
            try:
                response = self.client.embed(model=self.model, input=truncated_text)
                embeddings.append(response["embeddings"][0])
            except Exception as e:
                print(f"[DEBUG] Error embedding text {i}: {e}")
                raise
        return embeddings

    def embed_query(self, text: str) -> list[float]:
        truncated_text = self._truncate(text)
        response = self.client.embed(model=self.model, input=truncated_text)
        return response["embeddings"][0]

def is_valid_base64_image(content: str) -> bool:
    if len(content) < 500:
        return False
    if not re.match(r'^[A-Za-z0-9+/=]+$', content):
        return False
    try:
        decoded = b64decode(content, validate=True)
        return len(decoded) > 100
    except Exception:
        return False

def parse_docs(docs):
    """Split base64-encoded images and texts"""
    print(f"\n[DEBUG parse_docs] Received {len(docs)} documents")
    b64 = []
    text = []
    for i, doc in enumerate(docs):
        content = doc.page_content if hasattr(doc, 'page_content') else doc
        print(f"[DEBUG parse_docs] Doc {i}: type={type(doc).__name__}, len={len(content)}, preview={content[:80]}...")

        if is_valid_base64_image(content):
            b64.append(content)
            print(f"[DEBUG parse_docs] Doc {i}: Detected as BASE64 IMAGE")
        else:
            text.append(content)
            print(f"[DEBUG parse_docs] Doc {i}: Detected as TEXT")
    print(f"[DEBUG parse_docs] Result: {len(b64)} images, {len(text)} texts\n")
    return {"images": b64, "texts": text}


def build_prompt(kwargs):

    docs_by_type = kwargs["context"]
    user_question = kwargs["question"]

    print(f"\n[DEBUG build_prompt] Images: {len(docs_by_type['images'])}, Texts: {len(docs_by_type['texts'])}")
    
    context_text = ""
    if len(docs_by_type["texts"]) > 0:
        for text_content in docs_by_type["texts"]:
            context_text += text_content + "\n"

    prompt_template = f"""
    Answer the question based only on the following context, which can include text, tables, and the below image.
    Context: {context_text}
    Question: {user_question}
    """

    prompt_content = [{"type": "text", "text": prompt_template}]

    if len(docs_by_type["images"]) > 0:
        print(f"[DEBUG build_prompt] Adding {len(docs_by_type['images'])} images to prompt")
        for i, image in enumerate(docs_by_type["images"]):
            print(f"[DEBUG build_prompt] Image {i}: len={len(image)}")
            prompt_content.append(
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{image}"},
                }
            )
    else:
        print("[DEBUG build_prompt] No images found in context")

    print("Complete build prompt =======================>>>>>")
    return ChatPromptTemplate.from_messages(
        [
            HumanMessage(content=prompt_content),
        ]
    )
def summary_text(texts, tables):
    
    # Prompt
    prompt_text = """
    You are an assistant tasked with summarizing tables and text.
    Give a concise summary of the table or text.

    Respond only with the summary, no additionnal comment.
    Do not start your message by saying "Here is a summary" or anything like that.
    Just give the summary as it is.

    Table or text chunk: {element}

    """
    prompt = ChatPromptTemplate.from_template(prompt_text)

    # Summary chain
    # model = ChatGoogleGenerativeAI(
    #     model="gemini-2.0-flash-lite",
    #     temperature=0,
    #     max_tokens=65536,
    #     timeout=60,
    #     max_retries=2,
    #     api_key=GOOGLE_API_KEY,
    #     # other params...
    # )
    model = ChatOllama(model="llama3", temperature=0)
    # model = ChatGroq(temperature=0.5, model="llama-3.1-8b-instant", api_key=GROQ_API_KEY)
    summarize_chain = {"element": lambda x: x} | prompt | model | StrOutputParser()
    # Summarize text
    text_summaries = summarize_chain.batch(texts, {"max_concurrency": 3})

    tables_html = [table.metadata.text_as_html for table in tables]
    table_summaries = summarize_chain.batch(tables_html, {"max_concurrency": 3})
    return text_summaries, table_summaries

def summary_image(images):
    prompt_template = """Summarize this image concisely
    Focus on the key information: what it shows, main components, and purpose.
    Special case: If the photo contains a chart or graph, analyze it carefully and list all numbers, values, and data points shown in detail. Include every label, unit, axis name, category, and corresponding value. Describe the chart’s structure (for example, which values belong to which categories or time periods). For instance, if the image shows a rainfall chart for each month, list every month and the exact rainfall amount for each one. Present the results clearly and completely in a structured format such as a bullet list or table.
    """
    
    messages = [
        (
            "user",
            [
                {"type": "text", "text": prompt_template},
                {
                    "type": "image_url",
                    "image_url": {"url": "data:image/jpeg;base64,{image}"},
                },
            ],
        )
    ]

    prompt = ChatPromptTemplate.from_messages(messages)
    # print("start check image...")
    # for image in images:
    #     print(image)
    #     print("=========================")
    # print("end check image...")

    model = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0,
        max_tokens=65536,
        timeout=60,
        max_retries=2,
        api_key=GOOGLE_API_KEY,
    )
    chain = prompt | model | StrOutputParser()

    return chain.batch(images)
    
def filter_texts(chunks: list) -> list:
    return [
        chunk for chunk in chunks
        if isinstance(chunk, CompositeElement)
    ]

def filter_images_and_tables(chunks: list) -> tuple[list, list]:
    images, tables = [], []
    for chunk in chunks:
        if isinstance(chunk, Image):
            images.append(chunk)
        elif isinstance(chunk, Table):
            tables.append(chunk)
        elif isinstance(chunk, CompositeElement):
            for el in chunk.metadata.orig_elements or []:
                if isinstance(el, Image):
                    images.append(el)
                elif isinstance(el, Table):
                    tables.append(el)
    return images, tables

def filter_elements(chunks: list) -> dict:
    images, tables = filter_images_and_tables(chunks)
    return {
        "texts": filter_texts(chunks),
        "images": images,
        "tables": tables,
    }

def convert_img_to_base64(path: str) -> str:
    with open(path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def extract_images_from_pptx(pptx_path: str) -> list[str]:
    """Extract all images from PPTX file and return as base64 strings."""
    prs = PptxPresentation(pptx_path)
    images_base64 = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                images_base64.append(image_base64)
                print(f"[DEBUG] Slide {slide_num}: Found image, size={len(image_bytes)} bytes")
            
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = sub_shape.image
                        image_bytes = image.blob
                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        images_base64.append(image_base64)
                        print(f"[DEBUG] Slide {slide_num}: Found grouped image, size={len(image_bytes)} bytes")
    
    print(f"[DEBUG] Total images extracted: {len(images_base64)}")
    return images_base64

def extract_and_chunk_images(images_paths: list[str], chunk_size: int = 500, chunk_overlap: int = 50) -> dict:
    print(f"\n[DEBUG extract_and_chunk_images] Processing {len(images_paths)} images...")
    
    images = [convert_img_to_base64(path) for path in images_paths]
    print(f"[DEBUG] Converted {len(images)} images to base64")
    
    print("[DEBUG] Generating summaries...")
    summaries = summary_image(images)
    print(f"[DEBUG] Generated {len(summaries)} summaries")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ".", "*", " ", ""]
    )
    
    chunked_images = []
    chunked_summaries = []
    
    for img, summary in zip(images, summaries):
        chunks = text_splitter.split_text(summary)
        print(f"[DEBUG] Split summary into {len(chunks)} chunks")
        for i, chunk in enumerate(chunks):
            print(f"[DEBUG] Chunk {i}: {len(chunk)} chars")
            chunked_images.append(img)
            chunked_summaries.append(chunk)
    
    print(f"[DEBUG] Total: {len(images)} images -> {len(chunked_images)} chunks")
    
    return {
        "images": chunked_images,
        "summaries": chunked_summaries,
        "original_count": len(images),
        "chunk_count": len(chunked_images)
    }

def store_vector(store: InMemoryStore, data: dict, summaries: dict, collection_name: str = "multi_modal_rag"):
    print(f"\n{'='*50}")
    print("[DEBUG store_vector] Starting indexing...")
    print(f"{'='*50}")
    
    embeddings = OllamaEmbeddingsCustom(model="nomic-embed-text")
    
    chroma_host = os.getenv("CHROMA_HOST", "localhost")
    chroma_port = int(os.getenv("CHROMA_PORT", "8001"))
    client = chromadb.HttpClient(host=chroma_host, port=chroma_port)
    
    try:
        client.delete_collection(collection_name)
        print(f"[DEBUG store_vector] Deleted old collection: {collection_name}")
    except Exception:
        print(f"[DEBUG store_vector] No existing collection to delete: {collection_name}")
    
    vectorstore = Chroma(
        client=client,
        collection_name=collection_name,
        embedding_function=embeddings,
    )

    id_key = "doc_id"

    retriever = MultiVectorRetriever(
        vectorstore=vectorstore,
        docstore=store,
        id_key=id_key,
    )

    total_indexed = 0

    if data.get('image') and summaries.get('image'):
        img_ids = [str(uuid.uuid4()) for _ in data['image']]
        print(f"[DEBUG store_vector] Indexing {len(img_ids)} images...")
        
        summary_img = [
            Document(page_content=summary, metadata={id_key: img_ids[i]})
            for i, summary in enumerate(summaries['image'])
        ]
        retriever.vectorstore.add_documents(summary_img)
        retriever.docstore.mset(list(zip(img_ids, data['image'])))
        total_indexed += len(img_ids)
        
        print(f"[DEBUG store_vector] Added {len(img_ids)} images to vectorstore")

    if data.get('text') and summaries.get('text'):
        text_ids = [str(uuid.uuid4()) for _ in data['text']]
        print(f"[DEBUG store_vector] Indexing {len(text_ids)} texts...")
        
        summary_texts = [
            Document(page_content=summary, metadata={id_key: text_ids[i]})
            for i, summary in enumerate(summaries['text'])
        ]
        retriever.vectorstore.add_documents(summary_texts)
        retriever.docstore.mset(list(zip(text_ids, [str(t) for t in data['text']])))
        total_indexed += len(text_ids)
        
        print(f"[DEBUG store_vector] Added {len(text_ids)} texts to vectorstore")

    if data.get('table') and summaries.get('table'):
        table_ids = [str(uuid.uuid4()) for _ in data['table']]
        print(f"[DEBUG store_vector] Indexing {len(table_ids)} tables...")
        
        summary_tables = [
            Document(page_content=summary, metadata={id_key: table_ids[i]})
            for i, summary in enumerate(summaries['table'])
        ]
        retriever.vectorstore.add_documents(summary_tables)
        retriever.docstore.mset(list(zip(table_ids, [str(t) for t in data['table']])))
        total_indexed += len(table_ids)
        
        print(f"[DEBUG store_vector] Added {len(table_ids)} tables to vectorstore")

    print(f"\n[DEBUG store_vector] Verifying InMemoryStore...")
    all_keys = list(store.yield_keys())
    print(f"[DEBUG store_vector] InMemoryStore has {len(all_keys)} keys")
    
    for key in all_keys[:5]:
        data_item = store.mget([key])[0]
        if data_item:
            preview = data_item[:100] if isinstance(data_item, str) else str(data_item)[:100]
            print(f"[DEBUG store_vector] Key {key[:8]}...: len={len(data_item) if hasattr(data_item, '__len__') else 'N/A'}, preview={preview}...")

    print(f"\n[DEBUG store_vector] Indexing complete! Total: {total_indexed} items")
    print(f"{'='*50}\n")
    
    return {
        "total_indexed": total_indexed,
        "collection": collection_name,
        "store_keys": len(all_keys)
    }



def retrieve_data():
    embeddings = OllamaEmbeddingsCustom(model="nomic-embed-text")
    
    chroma_host = os.getenv("CHROMA_HOST", "localhost")
    chroma_port = int(os.getenv("CHROMA_PORT", "8001"))
    client = chromadb.HttpClient(host=chroma_host, port=chroma_port)
    
    vectorstore = Chroma(
        client=client,
        collection_name="multi_modal_rag",
        embedding_function=embeddings,
    )

    retriever = vectorstore.as_retriever(search_kwargs={"k": 1})

    # chain = (
    #     {
    #         "context": retriever | RunnableLambda(parse_docs),
    #         "question": RunnablePassthrough(),
    #     }
    #     | RunnableLambda(build_prompt)
    #     | ChatOllama(model="llama3", temperature=0)
    #     | StrOutputParser()
    # )

    # response = chain.invoke(
    #     "What is the attention mechanism?"
    # )

    # print(response)

    chain_with_sources = {
        "context": retriever | RunnableLambda(parse_docs),
        "question": RunnablePassthrough(),
    } | RunnablePassthrough().assign(
        response=(
            RunnableLambda(build_prompt)
            | ChatOllama(model="llava", temperature=0)
            | StrOutputParser()
        )
    )

    response = chain_with_sources.invoke(
        "What is Transformer model?"
    )

    print("Response:", response['response'])

    print("\n\nContext:")
    for text in response['context']['texts']:
        print(text.text)
        print("Page number: ", text.metadata.page_number)
        print("\n" + "-"*50 + "\n")
    for image in response['context']['images']:
        print(image)

def extract_images_from_docx(file_path: str) -> list[str]:
    """Extract images from DOCX file and return as base64 strings."""
    print(f"[DEBUG] Extracting images from DOCX: {file_path}")
    images_base64 = []
    
    output_dir = Path("./temp_images") / Path(file_path).stem
    output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        docx2txt.process(file_path, str(output_dir))
    except Exception as e:
        print(f"[DEBUG] Error processing DOCX: {e}")
        return []

    for image_path in output_dir.glob("*"):
        if image_path.is_file() and image_path.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            with open(image_path, "rb") as img_file:
                img_base64 = base64.b64encode(img_file.read()).decode('utf-8')
                images_base64.append(img_base64)
                print(f"[DEBUG] Found image: {image_path.name}, size={len(img_base64)} chars")
    
    print(f"[DEBUG] Total images extracted from DOCX: {len(images_base64)}")
    return images_base64

def main():

    output_path = "./content/"
    file_path = output_path + 'attention.pptx'

    # Reference: https://docs.unstructured.io/open-source/core-functionality/chunking
    # chunks = partition_pdf(
    #     filename=file_path,
    #     infer_table_structure=True,            # extract tables
    #     strategy="hi_res",                     # mandatory to infer tables

    #     extract_image_block_types=["Image", "Table"],   # Add 'Table' to list to extract image of tables
    #     # image_output_dir_path=output_path,   # if None, images and tables will saved in base64

    #     extract_image_block_to_payload=True,   # if true, will extract base64 for API usage

    #     chunking_strategy="by_title",          # or 'basic'
    #     max_characters=10000,                  # defaults to 500
    #     combine_text_under_n_chars=2000,       # defaults to 0
    #     new_after_n_chars=6000,

    #     # extract_images_in_pdf=True,          # deprecated
    # )

    
    chunks = partition(filename="./content/docs/test.docx")
    
    print(f"\n{'='*60}")
    print(f"Total chunks: {len(chunks)}")
    print(f"{'='*60}\n")
    print(chunks[0])
    for i, chunk in enumerate(chunks):
        chunk_type = type(chunk).__name__
        has_image = False
        image_base64 = None
        
        if isinstance(chunk, Image):
            has_image = True
            image_base64 = getattr(chunk.metadata, 'image_base64', None)
        
        if hasattr(chunk, 'metadata') and hasattr(chunk.metadata, 'orig_elements'):
            orig_elements = chunk.metadata.orig_elements or []
            for el in orig_elements:
                if isinstance(el, Image):
                    has_image = True
                    image_base64 = getattr(el.metadata, 'image_base64', None)
                    break
        
        text_preview = ""
        if hasattr(chunk, 'text') and chunk.text:
            text_preview = chunk.text[:60] + "..." if len(chunk.text) > 60 else chunk.text
        
        print(f"[{i}] {chunk_type}: image={has_image}, b64_len={len(image_base64) if image_base64 else 0}")
        if text_preview:
            print(f"    text: {text_preview}")
    
    elements = filter_elements(chunks)
    print(f"\n{'='*60}")
    print(f"partition_pptx result:")
    print(f"Texts: {len(elements['texts'])}, Images: {len(elements['images'])}, Tables: {len(elements['tables'])}")
    print(f"{'='*60}")
    
    print(f"\n{'='*60}")
    print("Using python-pptx to extract images directly:")
    print(f"{'='*60}")
    pptx_images = extract_images_from_pptx("./content/attension.pptx")
    print(f"Found {len(pptx_images)} images using python-pptx!")

    image_elements = extract_images_from_docx("./content/docs/test.docx")
    print(f"Extracted {len(image_elements)} images from DOCX.")
    
    with open("img.txt", "w") as f:
        for i, img_base64 in enumerate(image_elements):
            f.write(f"=== IMAGE {i+1} ===\n")
            f.write(img_base64)
            f.write("\n\n")
    print(f"Written {len(image_elements)} images to img.txt")

    # print(f"Texts: {len(elements['texts'])}")
    # print(f"Images: {len(elements['images'])}")
    # print(f"Tables: {len(elements['tables'])}")
    # images = [img.metadata.image_base64 for img in elements['images']]
    # print("START SUMMARY TEXT...")
    # t, tb = summary_text(elements['texts'],elements['tables'])
    # print("END SUMMARY TEXT...")
    # print("START EXTRACT AND CHUNK IMAGES...")
    # images_paths = ["./content/images/image.png"]
    
    # result = extract_and_chunk_images(images_paths, chunk_size=500, chunk_overlap=50)
    
    # images = result["images"]
    # sum_img = result["summaries"]
    
    # print(f"[DEBUG] Result: {result['original_count']} images -> {result['chunk_count']} chunks")
    # print("END EXTRACT AND CHUNK IMAGES...")
    # print("---------------------------")
    # embeddings = OllamaEmbeddingsCustom(model="nomic-embed-text")
    
    # chroma_host = os.getenv("CHROMA_HOST", "localhost")
    # chroma_port = int(os.getenv("CHROMA_PORT", "8001"))
    # client = chromadb.HttpClient(host=chroma_host, port=chroma_port)
    
    # collection_name = "multi_modal_rag"
    # try:
    #     client.delete_collection(collection_name)
    #     print(f"[DEBUG] Deleted old collection: {collection_name}")
    # except Exception:
    #     print(f"[DEBUG] No existing collection to delete")
    
    # vectorstore = Chroma(
    #     client=client,
    #     collection_name=collection_name,
    #     embedding_function=embeddings,
    # )
    
    # local_store = InMemoryStore()
    # id_key = "doc_id"
    
    # retriever = MultiVectorRetriever(
    #     vectorstore=vectorstore,
    #     docstore=local_store,
    #     id_key=id_key,
    # )
    
    # doc_ids = [str(uuid.uuid4()) for _ in sum_img]
    # print(f"[DEBUG] Generated {len(doc_ids)} doc_ids for {len(sum_img)} chunks")
    
    # summary_docs = [
    #     Document(page_content=summary, metadata={id_key: doc_ids[i]})
    #     for i, summary in enumerate(sum_img)
    # ]
    
    # print(f"[DEBUG] Adding {len(summary_docs)} documents to vectorstore...")
    # retriever.vectorstore.add_documents(summary_docs)
    
    # print(f"[DEBUG] Adding {len(images)} images to InMemoryStore...")
    # retriever.docstore.mset(list(zip(doc_ids, images)))
    
    # print("\n[DEBUG] Verifying stored data...")
    # stored_data = local_store.mget(doc_ids[:3])
    # for i, (doc_id, data) in enumerate(zip(doc_ids[:3], stored_data)):
    #     if data:
    #         preview = data[:100] if isinstance(data, str) else str(data)[:100]
    #         print(f"[DEBUG] ID {doc_id[:8]}...: len={len(data)}, preview={preview}...")
    
    # print("==========>> END STORE DATABASE <<")


if __name__ == "__main__":
    main()
